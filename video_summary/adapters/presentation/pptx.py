from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from typing import Optional

from video_summary.config import PipelineConfig
from video_summary.domain.models import ArtifactRecord, SceneSegment
from video_summary.services import format_ts


def add_textbox(slide: object, left: object, top: object, width: object, height: object, text: str, font_size: int = 18, bold: bool = False) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    text_frame = tb.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
    for inner_paragraph in text_frame.paragraphs:
        for run in inner_paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold


def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Optional[Path]:
    soffice = shutil.which("soffice")
    if not soffice:
        return None
    subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ],
        check=True,
    )
    pdf_path = output_dir / (pptx_path.stem + ".pdf")
    return pdf_path if pdf_path.exists() else None


class PptxPresentationGenerator:
    def generate(self, slides: list[SceneSegment], title: str, config: PipelineConfig) -> list[ArtifactRecord]:
        from pptx import Presentation
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt

        layout = config.layout()
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        for segment in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            header = (
                f"Segment {segment.index:03d} - {format_ts(segment.start)}-{format_ts(segment.end)} "
                f"- utterances: {segment.utterance_count}"
            )
            add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5), header, font_size=20, bold=True)

            slide.shapes.add_picture(segment.frame_path, Inches(0.4), Inches(0.8), width=Inches(7.8), height=Inches(5.9))

            tb = slide.shapes.add_textbox(Inches(8.4), Inches(0.85), Inches(4.3), Inches(5.85))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            lines = segment.text.splitlines() if segment.text else ["(no transcript text for this segment)"]
            for index, line in enumerate(lines):
                paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                paragraph.text = line
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(12)

        presentation.save(str(layout.slides_pptx))
        artifacts = [ArtifactRecord(name="slides_pptx", path=str(layout.slides_pptx), kind="presentation")]
        if config.export_pdf:
            pdf_path = convert_pptx_to_pdf(layout.slides_pptx, layout.output_dir)
            if pdf_path is not None:
                artifacts.append(ArtifactRecord(name="slides_pdf", path=str(pdf_path), kind="presentation"))
        return artifacts
