from __future__ import annotations

import argparse
import json
import math
import os
import shutil
import subprocess
import sys
import wave
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from scenedetect import detect
from scenedetect.detectors import AdaptiveDetector, ContentDetector, HashDetector
from tqdm import tqdm


@dataclass
class WordToken:
    start: float
    end: float
    text: str
    speaker: Optional[int] = None


@dataclass
class SpeakerTurn:
    start: float
    end: float
    speaker: int


@dataclass
class Utterance:
    start: float
    end: float
    speaker: int
    text: str


@dataclass
class SceneSegment:
    index: int
    start: float
    end: float
    frame_path: str
    text: str
    utterance_count: int


def run(cmd: Sequence[str], cwd: Optional[Path] = None, capture: bool = False) -> subprocess.CompletedProcess:
    kwargs = {
        "cwd": str(cwd) if cwd else None,
        "check": True,
        "text": True,
    }
    if capture:
        kwargs["stdout"] = subprocess.PIPE
        kwargs["stderr"] = subprocess.PIPE
    return subprocess.run(list(cmd), **kwargs)


def ffmpeg_supports_encoder(encoder: str) -> bool:
    try:
        out = run(["ffmpeg", "-hide_banner", "-encoders"], capture=True)
    except subprocess.CalledProcessError:
        return False
    return encoder in out.stdout


def resolve_ffmpeg_video_encoder(requested: str) -> str:
    if requested == "auto":
        return "h264_nvenc" if ffmpeg_supports_encoder("h264_nvenc") else "libx264"
    return requested


def video_encode_args(encoder: str) -> List[str]:
    if encoder == "h264_nvenc":
        return [
            "-c:v", "h264_nvenc",
            "-preset", "p5",
            "-cq", "23",
            "-b:v", "0",
            "-pix_fmt", "yuv420p",
        ]
    return [
        "-c:v", "libx264",
        "-preset", "medium",
        "-crf", "23",
        "-pix_fmt", "yuv420p",
    ]


def ensure_tool(name: str) -> None:
    if shutil.which(name) is None:
        raise RuntimeError(f"РќРµ РЅР°Р№РґРµРЅ РёРЅСЃС‚СЂСѓРјРµРЅС‚ '{name}' РІ PATH.")


def ffprobe_json(input_path: Path) -> Dict[str, Any]:
    cmd = [
        "ffprobe", "-v", "error",
        "-show_entries", "format=duration:stream=index,codec_type,width,height,r_frame_rate",
        "-of", "json", str(input_path),
    ]
    out = run(cmd, capture=True)
    return json.loads(out.stdout)


def parse_fps(r_frame_rate: str) -> float:
    if not r_frame_rate or r_frame_rate == "0/0":
        return 25.0
    num, den = r_frame_rate.split("/")
    num_f = float(num)
    den_f = float(den)
    return num_f / den_f if den_f else 25.0


def ensure_work_mp4(input_video: Path, work_video: Path, video_encoder: str) -> None:
    cmd = [
        "ffmpeg", "-y", "-i", str(input_video),
        "-map", "0:v:0",
        "-map", "0:a:0?",
        *video_encode_args(video_encoder),
        "-c:a", "aac",
        "-b:a", "128k",
        str(work_video),
    ]
    try:
        run(cmd)
    except subprocess.CalledProcessError:
        if video_encoder == "libx264":
            raise
        print(f"[WARN] ffmpeg video encoder '{video_encoder}' failed, retrying with libx264")
        fallback_cmd = [
            "ffmpeg", "-y", "-i", str(input_video),
            "-map", "0:v:0",
            "-map", "0:a:0?",
            *video_encode_args("libx264"),
            "-c:a", "aac",
            "-b:a", "128k",
            str(work_video),
        ]
        run(fallback_cmd)


def extract_audio(input_video: Path, audio_wav: Path) -> None:
    cmd = [
        "ffmpeg", "-y", "-i", str(input_video),
        "-vn",
        "-ac", "1",
        "-ar", "16000",
        "-c:a", "pcm_s16le",
        str(audio_wav),
    ]
    run(cmd)


def format_ts(seconds: float, srt: bool = False) -> str:
    seconds = max(0.0, float(seconds))
    ms = int(round((seconds - math.floor(seconds)) * 1000))
    whole = int(math.floor(seconds))
    s = whole % 60
    m = (whole // 60) % 60
    h = whole // 3600
    if srt:
        return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"
    cs = ms // 10
    return f"{h:d}:{m:02d}:{s:02d}.{cs:02d}"


def clean_text_spacing(text: str) -> str:
    text = " ".join(text.strip().split())
    text = text.replace(" ,", ",").replace(" .", ".").replace(" !", "!").replace(" ?", "?")
    text = text.replace(" :", ":").replace(" ;", ";")
    text = text.replace("( ", "(").replace(" )", ")")
    return text.strip()


def is_cuda_runtime_error(exc: BaseException) -> bool:
    message = str(exc).lower()
    cuda_markers = (
        "cublas",
        "cudnn",
        "cuda",
        "cannot be loaded",
        "not found",
        "failed to load",
    )
    return any(marker in message for marker in cuda_markers)


def normalize_compute_type(device: str, compute_type: str) -> str:
    if device == "cpu" and compute_type in {"float16", "int8_float16", "bfloat16"}:
        return "int8"
    return compute_type


def transcribe_audio_faster_whisper(
    audio_path: Path,
    model_name: str,
    language: Optional[str],
    device: str,
    compute_type: str,
    verbose: bool = True,
) -> Tuple[List[WordToken], Dict[str, Any]]:
    from faster_whisper import WhisperModel

    resolved_device = device
    resolved_compute_type = normalize_compute_type(device, compute_type)
    if resolved_compute_type != compute_type and verbose:
        print(f"[ASR] compute_type={compute_type} is not suitable for CPU, using {resolved_compute_type}")

    def run_transcribe(device_name: str, compute_name: str) -> Tuple[List[Any], Any]:
        if verbose:
            print(f"[ASR] Loading faster-whisper model={model_name} device={device_name} compute_type={compute_name}")
        model = WhisperModel(model_name, device=device_name, compute_type=compute_name)
        segments_gen, info_obj = model.transcribe(
            str(audio_path),
            beam_size=5,
            language=language,
            word_timestamps=True,
            vad_filter=True,
            vad_parameters=dict(min_silence_duration_ms=500),
            condition_on_previous_text=False,
        )
        return list(segments_gen), info_obj

    if verbose:
        print(f"[ASR] Requested device={device} compute_type={compute_type}")

    try:
        segments, info = run_transcribe(resolved_device, resolved_compute_type)
    except RuntimeError as exc:
        if resolved_device == "cuda" and is_cuda_runtime_error(exc):
            resolved_device = "cpu"
            resolved_compute_type = "int8"
            print(f"[WARN] CUDA ASR failed ({exc}). Falling back to CPU ({resolved_compute_type}).")
            segments, info = run_transcribe(resolved_device, resolved_compute_type)
        else:
            raise

    words: List[WordToken] = []
    for seg in segments:
        seg_words = getattr(seg, "words", None)
        if seg_words:
            for w in seg_words:
                if w.start is None or w.end is None:
                    continue
                token_text = clean_text_spacing(getattr(w, "word", ""))
                if token_text:
                    words.append(WordToken(float(w.start), float(w.end), token_text))
        else:
            text = clean_text_spacing(getattr(seg, "text", ""))
            if text:
                words.append(WordToken(float(seg.start), float(seg.end), text))

    meta = {
        "language": getattr(info, "language", language),
        "language_probability": getattr(info, "language_probability", None),
        "model_name": model_name,
        "device": resolved_device,
        "compute_type": resolved_compute_type,
        "segment_count": len(segments),
        "word_count": len(words),
    }
    return words, meta


def diarize_audio_pyannote(
    audio_path: Path,
    hf_token: str,
    device: str,
    num_speakers: Optional[int] = None,
    min_speakers: Optional[int] = None,
    max_speakers: Optional[int] = None,
    verbose: bool = True,
) -> List[SpeakerTurn]:
    import torch
    from pyannote.audio import Pipeline

    if verbose:
        print("[DIAR] Loading pyannote community-1 diarization pipeline")

    pipeline = Pipeline.from_pretrained(
        "pyannote/speaker-diarization-community-1",
        token=hf_token,
    )
    if device == "cuda":
        try:
            pipeline.to(torch.device("cuda"))
        except RuntimeError as exc:
            if is_cuda_runtime_error(exc):
                print(f"[WARN] CUDA diarization failed ({exc}). Falling back to CPU.")
            else:
                raise

    kwargs: Dict[str, Any] = {}
    if num_speakers is not None:
        kwargs["num_speakers"] = num_speakers
    if min_speakers is not None:
        kwargs["min_speakers"] = min_speakers
    if max_speakers is not None:
        kwargs["max_speakers"] = max_speakers

    audio_input: Any = str(audio_path)
    try:
        with wave.open(str(audio_path), "rb") as wav_file:
            channels = wav_file.getnchannels()
            sample_width = wav_file.getsampwidth()
            sample_rate = wav_file.getframerate()
            frame_count = wav_file.getnframes()
            raw_bytes = wav_file.readframes(frame_count)

        if sample_width == 2:
            pcm = np.frombuffer(raw_bytes, dtype=np.int16).astype(np.float32) / 32768.0
            if channels > 1:
                pcm = pcm.reshape(-1, channels).T
            else:
                pcm = pcm.reshape(1, -1)
            waveform = torch.from_numpy(np.ascontiguousarray(pcm))
            audio_input = {"waveform": waveform, "sample_rate": sample_rate, "uri": audio_path.stem}
            if verbose:
                print("[DIAR] Using in-memory waveform input to avoid AudioDecoder dependency.")
        elif verbose:
            print(f"[DIAR] WAV sample width {sample_width} bytes is unsupported for in-memory fallback; using file path.")
    except (wave.Error, OSError, ValueError) as exc:
        if verbose:
            print(f"[DIAR] Could not preload audio ({exc}); using file path for diarization.")

    diar_output = pipeline(audio_input, **kwargs)
    ann = getattr(diar_output, "speaker_diarization", diar_output)

    raw_turns: List[Tuple[float, float, str]] = []
    if hasattr(ann, "itertracks"):
        for turn, _track, label in ann.itertracks(yield_label=True):
            raw_turns.append((float(turn.start), float(turn.end), str(label)))
    else:
        # Fallback for alternate API shapes
        for item in ann:
            if len(item) == 2:
                turn, label = item
                raw_turns.append((float(turn.start), float(turn.end), str(label)))
            elif len(item) == 3:
                turn, _track, label = item
                raw_turns.append((float(turn.start), float(turn.end), str(label)))

    raw_turns.sort(key=lambda x: (x[0], x[1]))

    label_map: Dict[str, int] = {}
    next_id = 1
    out: List[SpeakerTurn] = []
    for start, end, label in raw_turns:
        if label not in label_map:
            label_map[label] = next_id
            next_id += 1
        out.append(SpeakerTurn(start=start, end=end, speaker=label_map[label]))

    return out


def interval_overlap(a_start: float, a_end: float, b_start: float, b_end: float) -> float:
    return max(0.0, min(a_end, b_end) - max(a_start, b_start))


def assign_speakers_to_words(words: List[WordToken], turns: List[SpeakerTurn]) -> List[WordToken]:
    if not turns:
        for w in words:
            w.speaker = 1
        return words

    t_idx = 0
    n = len(turns)

    for w in words:
        best_speaker = None
        best_overlap = -1.0
        midpoint = (w.start + w.end) / 2.0

        while t_idx + 1 < n and turns[t_idx + 1].end <= w.start:
            t_idx += 1

        candidate_indices = range(max(0, t_idx - 2), min(n, t_idx + 4))
        for idx in candidate_indices:
            t = turns[idx]
            ov = interval_overlap(w.start, w.end, t.start, t.end)
            if ov > best_overlap:
                best_overlap = ov
                best_speaker = t.speaker

        if best_overlap <= 0:
            nearest = min(
                turns,
                key=lambda t: abs(((t.start + t.end) / 2.0) - midpoint),
            )
            best_speaker = nearest.speaker

        w.speaker = best_speaker if best_speaker is not None else 1

    return words


def build_utterances(words: List[WordToken], gap_sec: float = 0.8) -> List[Utterance]:
    if not words:
        return []

    words = sorted(words, key=lambda x: (x.start, x.end))
    out: List[Utterance] = []

    cur_words = [words[0].text]
    cur_start = words[0].start
    cur_end = words[0].end
    cur_speaker = words[0].speaker or 1

    for w in words[1:]:
        speaker = w.speaker or 1
        gap = w.start - cur_end
        if speaker != cur_speaker or gap > gap_sec:
            out.append(Utterance(
                start=cur_start,
                end=cur_end,
                speaker=cur_speaker,
                text=clean_text_spacing(" ".join(cur_words)),
            ))
            cur_words = [w.text]
            cur_start = w.start
            cur_end = w.end
            cur_speaker = speaker
        else:
            cur_words.append(w.text)
            cur_end = w.end

    out.append(Utterance(
        start=cur_start,
        end=cur_end,
        speaker=cur_speaker,
        text=clean_text_spacing(" ".join(cur_words)),
    ))
    return out


def build_subtitle_chunks(
    words: List[WordToken],
    max_chars: int = 84,
    max_duration: float = 4.5,
    gap_sec: float = 0.7,
) -> List[Utterance]:
    if not words:
        return []

    words = sorted(words, key=lambda x: (x.start, x.end))
    out: List[Utterance] = []

    current: List[WordToken] = [words[0]]
    for w in words[1:]:
        cur_text = clean_text_spacing(" ".join(x.text for x in current))
        cur_duration = current[-1].end - current[0].start
        same_speaker = (w.speaker or 1) == (current[-1].speaker or 1)
        gap = w.start - current[-1].end
        next_text = clean_text_spacing(cur_text + " " + w.text)

        should_split = (
            not same_speaker or
            gap > gap_sec or
            cur_duration >= max_duration or
            len(next_text) > max_chars
        )
        if should_split:
            speaker = current[0].speaker or 1
            text = clean_text_spacing(" ".join(x.text for x in current))
            out.append(Utterance(current[0].start, current[-1].end, speaker, f"[{speaker}] {text}"))
            current = [w]
        else:
            current.append(w)

    if current:
        speaker = current[0].speaker or 1
        text = clean_text_spacing(" ".join(x.text for x in current))
        out.append(Utterance(current[0].start, current[-1].end, speaker, f"[{speaker}] {text}"))

    return out


def write_srt(subtitles: List[Utterance], out_path: Path) -> None:
    lines: List[str] = []
    for i, sub in enumerate(subtitles, start=1):
        lines.append(str(i))
        lines.append(f"{format_ts(sub.start, srt=True)} --> {format_ts(sub.end, srt=True)}")
        lines.append(sub.text)
        lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def ass_escape(text: str) -> str:
    return text.replace("\\", r"\\").replace("{", r"\{").replace("}", r"\}").replace("\n", r"\N")


def write_ass(subtitles: List[Utterance], out_path: Path, width: int, height: int) -> None:
    header = f"""[Script Info]
ScriptType: v4.00+
PlayResX: {width}
PlayResY: {height}
ScaledBorderAndShadow: yes
WrapStyle: 0
YCbCr Matrix: TV.601

[V4+ Styles]
Format: Name,Fontname,Fontsize,PrimaryColour,SecondaryColour,OutlineColour,BackColour,Bold,Italic,Underline,StrikeOut,ScaleX,ScaleY,Spacing,Angle,BorderStyle,Outline,Shadow,Alignment,MarginL,MarginR,MarginV,Encoding
Style: Default,Arial,28,&H00FFFFFF,&H0000FFFF,&H00101010,&H80000000,0,0,0,0,100,100,0,0,1,2,0,2,40,40,26,1

[Events]
Format: Layer,Start,End,Style,Name,MarginL,MarginR,MarginV,Effect,Text
"""
    lines = [header.rstrip()]
    for sub in subtitles:
        text = ass_escape(sub.text)
        lines.append(
            f"Dialogue: 0,{format_ts(sub.start)},{format_ts(sub.end)},Default,,0,0,0,,{text}"
        )
    out_path.write_text("\n".join(lines), encoding="utf-8")


def burn_subtitles_hard(work_video: Path, ass_file: Path, output_video: Path, video_encoder: str) -> None:
    # Use relative subtitle filename from its own directory to avoid Windows path escaping issues.
    cmd = [
        "ffmpeg", "-y",
        "-i", str(work_video),
        "-vf", f"subtitles={ass_file.name}",
        *video_encode_args(video_encoder),
        "-c:a", "aac",
        "-b:a", "128k",
        str(output_video),
    ]
    try:
        run(cmd, cwd=ass_file.parent)
    except subprocess.CalledProcessError:
        if video_encoder == "libx264":
            raise
        print(f"[WARN] ffmpeg video encoder '{video_encoder}' failed for hard subtitles, retrying with libx264")
        fallback_cmd = [
            "ffmpeg", "-y",
            "-i", str(work_video),
            "-vf", f"subtitles={ass_file.name}",
            *video_encode_args("libx264"),
            "-c:a", "aac",
            "-b:a", "128k",
            str(output_video),
        ]
        run(fallback_cmd, cwd=ass_file.parent)


def mux_subtitles_soft(work_video: Path, srt_file: Path, output_video: Path) -> None:
    cmd = [
        "ffmpeg", "-y",
        "-i", str(work_video),
        "-i", str(srt_file),
        "-map", "0:v:0",
        "-map", "0:a:0?",
        "-map", "1:0",
        "-c:v", "copy",
        "-c:a", "copy",
        "-c:s", "mov_text",
        str(output_video),
    ]
    run(cmd)


def extract_frame(video_path: Path, ts_sec: float, out_path: Path) -> None:
    cmd = [
        "ffmpeg", "-y",
        "-ss", f"{ts_sec:.3f}",
        "-i", str(video_path),
        "-frames:v", "1",
        "-q:v", "2",
        str(out_path),
    ]
    run(cmd)


def detect_scene_boundaries(
    video_path: Path,
    fps: float,
    detector_name: str,
    threshold: Optional[float],
    min_scene_sec: float,
) -> List[Tuple[float, float]]:
    min_scene_len = max(15, int(round(fps * min_scene_sec)))

    if detector_name == "content":
        detector = ContentDetector(threshold=threshold if threshold is not None else 27.0, min_scene_len=min_scene_len)
    elif detector_name == "adaptive":
        detector = AdaptiveDetector(adaptive_threshold=threshold if threshold is not None else 3.0, min_scene_len=min_scene_len)
    elif detector_name == "hash":
        detector = HashDetector(threshold=threshold if threshold is not None else 0.18, min_scene_len=min_scene_len)
    else:
        raise ValueError(f"Unknown scene detector: {detector_name}")

    scene_list = detect(str(video_path), detector, show_progress=True)
    scenes: List[Tuple[float, float]] = []
    for start_tc, end_tc in scene_list:
        scenes.append((float(start_tc.get_seconds()), float(end_tc.get_seconds())))
    return scenes


def merge_short_scenes(scenes: List[Tuple[float, float]], min_keep_sec: float) -> List[Tuple[float, float]]:
    if not scenes:
        return scenes
    merged: List[Tuple[float, float]] = []
    cur_start, cur_end = scenes[0]
    for start, end in scenes[1:]:
        if (cur_end - cur_start) < min_keep_sec:
            cur_end = end
        else:
            merged.append((cur_start, cur_end))
            cur_start, cur_end = start, end
    merged.append((cur_start, cur_end))
    return merged


def decide_has_presentation(
    scenes: List[Tuple[float, float]],
    total_duration: float,
) -> bool:
    if len(scenes) < 2:
        return False
    long_scenes = [s for s in scenes if (s[1] - s[0]) >= 8.0]
    long_coverage = sum(end - start for start, end in long_scenes) / max(1e-6, total_duration)
    return len(long_scenes) >= 2 and long_coverage >= 0.35


def utterances_overlapping(
    utterances: List[Utterance],
    start: float,
    end: float,
    pad: float = 1.0,
) -> List[Utterance]:
    qs = max(0.0, start - pad)
    qe = end + pad
    out = []
    for u in utterances:
        if interval_overlap(u.start, u.end, qs, qe) > 0:
            out.append(u)
    return out


def trim_slide_text(text: str, max_chars: int = 1700) -> str:
    text = text.strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "вЂ¦"


def build_slide_deck(
    scenes: List[Tuple[float, float]],
    utterances: List[Utterance],
    work_video: Path,
    out_dir: Path,
    force_single_slide: bool,
) -> List[SceneSegment]:
    frames_dir = out_dir / "frames"
    frames_dir.mkdir(parents=True, exist_ok=True)

    slides: List[SceneSegment] = []

    if force_single_slide or not scenes:
        probe = ffprobe_json(work_video)
        duration = float(probe["format"]["duration"])
        ts = duration / 2.0
        frame = frames_dir / "scene_001.jpg"
        extract_frame(work_video, ts, frame)
        text = "\n".join(f"[{u.speaker}] {u.text}" for u in utterances)
        slides.append(SceneSegment(
            index=1,
            start=0.0,
            end=duration,
            frame_path=str(frame),
            text=trim_slide_text(text),
            utterance_count=len(utterances),
        ))
        return slides

    for idx, (start, end) in enumerate(scenes, start=1):
        ts = max(start, min((start + end) / 2.0, max(start + 0.1, end - 0.1)))
        frame = frames_dir / f"scene_{idx:03d}.jpg"
        extract_frame(work_video, ts, frame)
        utts = utterances_overlapping(utterances, start, end, pad=1.0)
        text = "\n".join(f"[{u.speaker}] {u.text}" for u in utts)
        slides.append(SceneSegment(
            index=idx,
            start=start,
            end=end,
            frame_path=str(frame),
            text=trim_slide_text(text),
            utterance_count=len(utts),
        ))
    return slides


def add_textbox(slide, left, top, width, height, text: str, font_size: int = 18, bold: bool = False) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold


def create_pptx(slides_data: List[SceneSegment], pptx_path: Path, title: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for seg in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        header = f"РЎРµРіРјРµРЅС‚ {seg.index:03d} вЂў {format_ts(seg.start, srt=False)}вЂ“{format_ts(seg.end, srt=False)} вЂў СЂРµРїР»РёРє: {seg.utterance_count}"
        add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5), header, font_size=20, bold=True)

        img_left = Inches(0.4)
        img_top = Inches(0.8)
        img_width = Inches(7.8)
        img_height = Inches(5.9)

        slide.shapes.add_picture(seg.frame_path, img_left, img_top, width=img_width, height=img_height)

        notes_left = Inches(8.4)
        notes_top = Inches(0.85)
        notes_width = Inches(4.3)
        notes_height = Inches(5.85)

        tb = slide.shapes.add_textbox(notes_left, notes_top, notes_width, notes_height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for i, line in enumerate(seg.text.splitlines() if seg.text else ["(РЅРµС‚ С‚РµРєСЃС‚Р° РґР»СЏ СЌС‚РѕРіРѕ СЃРµРіРјРµРЅС‚Р°)"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(12)

    prs.save(str(pptx_path))


def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Optional[Path]:
    soffice = shutil.which("soffice")
    if not soffice:
        return None
    cmd = [
        soffice,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]
    run(cmd)
    pdf_path = output_dir / (pptx_path.stem + ".pdf")
    return pdf_path if pdf_path.exists() else None


def write_transcript_txt(utterances: List[Utterance], out_path: Path) -> None:
    lines: List[str] = []
    for u in utterances:
        lines.append(f"[{format_ts(u.start, srt=False)} -> {format_ts(u.end, srt=False)}] [S{u.speaker}]")
        lines.append(u.text)
        lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def save_json(obj: Any, out_path: Path) -> None:
    def default(o: Any) -> Any:
        if hasattr(o, "__dict__"):
            return o.__dict__
        if hasattr(o, "_asdict"):
            return o._asdict()
        raise TypeError(f"Cannot serialize type {type(o)}")
    out_path.write_text(json.dumps(obj, ensure_ascii=False, indent=2, default=default), encoding="utf-8")


STEP_ORDER = (
    "prepare",
    "asr",
    "diarize",
    "align",
    "scenes",
    "slides",
    "write",
    "render",
)
STEP_NUMBERS = {name: idx for idx, name in enumerate(STEP_ORDER, start=1)}


def load_json(in_path: Path) -> Any:
    return json.loads(in_path.read_text(encoding="utf-8"))


def state_json_path(out_dir: Path) -> Path:
    return out_dir / "pipeline_state.json"


def word_from_dict(data: Dict[str, Any]) -> WordToken:
    return WordToken(
        start=float(data["start"]),
        end=float(data["end"]),
        text=str(data["text"]),
        speaker=int(data["speaker"]) if data.get("speaker") is not None else None,
    )


def speaker_turn_from_dict(data: Dict[str, Any]) -> SpeakerTurn:
    return SpeakerTurn(
        start=float(data["start"]),
        end=float(data["end"]),
        speaker=int(data["speaker"]),
    )


def utterance_from_dict(data: Dict[str, Any]) -> Utterance:
    return Utterance(
        start=float(data["start"]),
        end=float(data["end"]),
        speaker=int(data["speaker"]),
        text=str(data["text"]),
    )


def scene_segment_from_dict(data: Dict[str, Any]) -> SceneSegment:
    return SceneSegment(
        index=int(data["index"]),
        start=float(data["start"]),
        end=float(data["end"]),
        frame_path=str(data["frame_path"]),
        text=str(data["text"]),
        utterance_count=int(data["utterance_count"]),
    )


def load_pipeline_state(state_path: Path) -> Dict[str, Any]:
    if not state_path.exists():
        raise FileNotFoundError(
            f"РќРµ РЅР°Р№РґРµРЅ С„Р°Р№Р» СЃРѕСЃС‚РѕСЏРЅРёСЏ '{state_path}'. "
            "Р”Р»СЏ Р·Р°РїСѓСЃРєР° СЃ РїСЂРѕРјРµР¶СѓС‚РѕС‡РЅРѕРіРѕ С€Р°РіР° СЃРЅР°С‡Р°Р»Р° РІС‹РїРѕР»РЅРё РїРѕР»РЅС‹Р№ РїСЂРѕРіРѕРЅ СЃ `--keep-work-files`."
        )
    data = load_json(state_path)
    if not isinstance(data, dict):
        raise RuntimeError(f"РќРµРєРѕСЂСЂРµРєС‚РЅС‹Р№ С„РѕСЂРјР°С‚ state-С„Р°Р№Р»Р°: {state_path}")
    return data


def save_pipeline_state(state_path: Path, state: Dict[str, Any]) -> None:
    state = dict(state)
    state["state_version"] = 1
    save_json(state, state_path)


def require_existing_file(path: Path, description: str, hint: Optional[str] = None) -> None:
    if path.exists():
        return
    msg = f"РќРµ РЅР°Р№РґРµРЅ {description}: {path}"
    if hint:
        msg += f". {hint}"
    raise FileNotFoundError(msg)


def validate_resume_input(state: Dict[str, Any], input_video: Path, start_from: str) -> None:
    if start_from == "prepare":
        return
    saved_input = state.get("input_video")
    if saved_input and Path(saved_input) != input_video:
        raise RuntimeError(
            "РќРµР»СЊР·СЏ РїСЂРѕРґРѕР»Р¶РёС‚СЊ РїР°Р№РїР»Р°Р№РЅ СЃ РґСЂСѓРіРёРј РІС…РѕРґРЅС‹Рј РІРёРґРµРѕ. "
            f"state-С„Р°Р№Р» СЃРѕР±СЂР°РЅ РґР»СЏ '{saved_input}', Р° РїРµСЂРµРґР°РЅ '{input_video}'."
        )


def load_prepare_artifacts(
    state: Dict[str, Any],
    work_video: Path,
    audio_wav: Path,
) -> Tuple[float, int, int, float]:
    require_existing_file(
        work_video,
        "РїРѕРґРіРѕС‚РѕРІР»РµРЅРЅС‹Р№ РІРёРґРµРѕС„Р°Р№Р» work.mp4",
        "Р—Р°РїСѓСЃС‚Рё СЃРЅР°С‡Р°Р»Р° СЃ `--start-from prepare --keep-work-files`.",
    )
    require_existing_file(
        audio_wav,
        "РїРѕРґРіРѕС‚РѕРІР»РµРЅРЅС‹Р№ Р°СѓРґРёРѕС„Р°Р№Р» audio.wav",
        "Р—Р°РїСѓСЃС‚Рё СЃРЅР°С‡Р°Р»Р° СЃ `--start-from prepare --keep-work-files`.",
    )
    video_meta = state.get("video")
    if not isinstance(video_meta, dict):
        raise RuntimeError("Р’ pipeline_state.json РЅРµС‚ СЃРµРєС†РёРё `video` РїРѕСЃР»Рµ С€Р°РіР° prepare.")
    return (
        float(state["duration_sec"]),
        int(video_meta["width"]),
        int(video_meta["height"]),
        float(video_meta["fps"]),
    )


def load_asr_artifacts(state: Dict[str, Any]) -> Tuple[List[WordToken], Dict[str, Any]]:
    raw_words = state.get("asr_words")
    asr_meta = state.get("asr_meta")
    if not isinstance(raw_words, list) or not isinstance(asr_meta, dict):
        raise RuntimeError("Р’ pipeline_state.json РЅРµС‚ Р°СЂС‚РµС„Р°РєС‚РѕРІ С€Р°РіР° asr (`asr_words`, `asr_meta`).")
    return [word_from_dict(item) for item in raw_words], asr_meta


def load_diarization_artifacts(state: Dict[str, Any]) -> List[SpeakerTurn]:
    raw_turns = state.get("speaker_turns")
    if not isinstance(raw_turns, list):
        raise RuntimeError("Р’ pipeline_state.json РЅРµС‚ `speaker_turns` РїРѕСЃР»Рµ С€Р°РіР° diarize.")
    return [speaker_turn_from_dict(item) for item in raw_turns]


def load_alignment_artifacts(state: Dict[str, Any]) -> Tuple[List[WordToken], List[Utterance], List[Utterance]]:
    raw_words = state.get("words")
    raw_utterances = state.get("utterances")
    raw_subtitles = state.get("subtitles")
    if not isinstance(raw_words, list) or not isinstance(raw_utterances, list) or not isinstance(raw_subtitles, list):
        raise RuntimeError("Р’ pipeline_state.json РЅРµС‚ РґР°РЅРЅС‹С… С€Р°РіР° align (`words`, `utterances`, `subtitles`).")
    return (
        [word_from_dict(item) for item in raw_words],
        [utterance_from_dict(item) for item in raw_utterances],
        [utterance_from_dict(item) for item in raw_subtitles],
    )


def load_scene_artifacts(state: Dict[str, Any]) -> Tuple[List[Tuple[float, float]], bool]:
    raw_scenes = state.get("scenes")
    if not isinstance(raw_scenes, list) or "has_presentation" not in state:
        raise RuntimeError("Р’ pipeline_state.json РЅРµС‚ Р°СЂС‚РµС„Р°РєС‚РѕРІ С€Р°РіР° scenes (`scenes`, `has_presentation`).")
    return [(float(start), float(end)) for start, end in raw_scenes], bool(state["has_presentation"])


def load_slide_artifacts(state: Dict[str, Any], pptx_path: Path, export_pdf: bool, pdf_path: Path) -> List[SceneSegment]:
    raw_segments = state.get("slide_segments")
    if not isinstance(raw_segments, list):
        raise RuntimeError("Р’ pipeline_state.json РЅРµС‚ `slide_segments` РїРѕСЃР»Рµ С€Р°РіР° slides.")
    require_existing_file(pptx_path, "СЃРіРµРЅРµСЂРёСЂРѕРІР°РЅРЅС‹Р№ slides.pptx")
    if export_pdf:
        require_existing_file(pdf_path, "СЃРіРµРЅРµСЂРёСЂРѕРІР°РЅРЅС‹Р№ slides.pdf")
    return [scene_segment_from_dict(item) for item in raw_segments]


def step_enabled(start_from: str, step_name: str) -> bool:
    return STEP_NUMBERS[start_from] <= STEP_NUMBERS[step_name]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Local pipeline for meeting video transcription + diarization + slides")
    parser.add_argument("--input", required=True, help="Input video file (.webm, .mp4, ...)")
    parser.add_argument("--output-dir", required=True, help="Output directory")
    parser.add_argument("--hf-token", default=os.environ.get("HF_TOKEN"), help="Hugging Face token for pyannote")
    parser.add_argument("--language", default=None, help="Language code for ASR, e.g. ru, en")
    parser.add_argument("--model", default="large-v3", help="faster-whisper model: large-v3 | distil-large-v3 | turbo")
    parser.add_argument("--device", default="cuda", choices=["cuda", "cpu"], help="Inference device")
    parser.add_argument("--compute-type", default="float16", help="faster-whisper compute type")
    parser.add_argument(
        "--ffmpeg-video-encoder",
        default="auto",
        choices=["auto", "h264_nvenc", "libx264"],
        help="Video encoder for ffmpeg. 'auto' prefers NVIDIA NVENC and falls back to libx264.",
    )
    parser.add_argument("--presentation", default="auto", choices=["auto", "yes", "no"], help="How to build slides")
    parser.add_argument("--scene-detector", default="content", choices=["content", "adaptive", "hash"])
    parser.add_argument("--scene-threshold", type=float, default=None)
    parser.add_argument("--min-scene-sec", type=float, default=5.0)
    parser.add_argument("--num-speakers", type=int, default=None)
    parser.add_argument("--min-speakers", type=int, default=None)
    parser.add_argument("--max-speakers", type=int, default=None)
    parser.add_argument("--subtitle-max-chars", type=int, default=84)
    parser.add_argument("--subtitle-max-duration", type=float, default=4.5)
    parser.add_argument("--transcript-gap", type=float, default=0.8)
    parser.add_argument(
        "--start-from",
        default="prepare",
        choices=list(STEP_ORDER),
        help="Start pipeline from a specific step. Useful for debugging with existing intermediate artifacts.",
    )
    parser.add_argument("--export-pdf", action="store_true")
    parser.add_argument("--keep-work-files", action="store_true")
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    input_video = Path(args.input).expanduser().resolve()
    out_dir = Path(args.output_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    work_dir = out_dir / "_work"
    work_dir.mkdir(parents=True, exist_ok=True)
    state_path = state_json_path(out_dir)
    work_video = work_dir / "work.mp4"
    audio_wav = work_dir / "audio.wav"
    pptx_path = out_dir / "slides.pptx"
    pdf_path = out_dir / "slides.pdf"
    transcript_txt = out_dir / "transcript.txt"
    transcript_json = out_dir / "transcript.json"
    srt_path = out_dir / "subtitles.srt"
    ass_path = out_dir / "subtitles.ass"
    hard_mp4 = out_dir / "video_subtitled.mp4"
    soft_mp4 = out_dir / "video_softsubs.mp4"

    if not input_video.exists():
        raise FileNotFoundError(input_video)

    ensure_tool("ffmpeg")
    ensure_tool("ffprobe")

    state: Dict[str, Any] = {}
    if args.start_from != "prepare":
        state = load_pipeline_state(state_path)
    validate_resume_input(state, input_video, args.start_from)

    if step_enabled(args.start_from, "diarize") and not args.hf_token:
        raise RuntimeError("РќСѓР¶РµРЅ Hugging Face token РґР»СЏ pyannote (`--hf-token` РёР»Рё env `HF_TOKEN`).")

    video_encoder = resolve_ffmpeg_video_encoder(args.ffmpeg_video_encoder)
    print(f"[INFO] ffmpeg video encoder: {video_encoder}")
    state["input_video"] = str(input_video)
    state["start_from"] = args.start_from

    if step_enabled(args.start_from, "prepare"):
        print("[1/8] Preparing media")
        ensure_work_mp4(input_video, work_video, video_encoder=video_encoder)
        extract_audio(work_video, audio_wav)

        probe = ffprobe_json(work_video)
        duration = float(probe["format"]["duration"])
        vstream = next((s for s in probe.get("streams", []) if s.get("codec_type") == "video"), None)
        width = int(vstream.get("width", 1920)) if vstream else 1920
        height = int(vstream.get("height", 1080)) if vstream else 1080
        fps = parse_fps(vstream.get("r_frame_rate", "25/1")) if vstream else 25.0
        state.update({
            "duration_sec": duration,
            "video": {"width": width, "height": height, "fps": fps},
        })
        save_pipeline_state(state_path, state)
    else:
        print("[1/8] Skipping media preparation, reusing existing artifacts")
        duration, width, height, fps = load_prepare_artifacts(state, work_video, audio_wav)

    if step_enabled(args.start_from, "asr"):
        print("[2/8] Running ASR")
        asr_words, asr_meta = transcribe_audio_faster_whisper(
            audio_wav,
            model_name=args.model,
            language=args.language,
            device=args.device,
            compute_type=args.compute_type,
        )
        state.update({
            "asr_meta": asr_meta,
            "asr_words": [asdict(x) for x in asr_words],
        })
        save_pipeline_state(state_path, state)
    else:
        print("[2/8] Skipping ASR, loading cached results")
        asr_words, asr_meta = load_asr_artifacts(state)

    if step_enabled(args.start_from, "diarize"):
        print("[3/8] Running speaker diarization")
        turns = diarize_audio_pyannote(
            audio_wav,
            hf_token=args.hf_token,
            device=args.device,
            num_speakers=args.num_speakers,
            min_speakers=args.min_speakers,
            max_speakers=args.max_speakers,
        )
        state["speaker_turns"] = [asdict(x) for x in turns]
        save_pipeline_state(state_path, state)
    else:
        print("[3/8] Skipping diarization, loading cached results")
        turns = load_diarization_artifacts(state)

    if step_enabled(args.start_from, "align"):
        print("[4/8] Aligning speakers with words")
        words = assign_speakers_to_words(asr_words, turns)
        utterances = build_utterances(words, gap_sec=args.transcript_gap)
        subtitles = build_subtitle_chunks(
            words,
            max_chars=args.subtitle_max_chars,
            max_duration=args.subtitle_max_duration,
            gap_sec=min(0.7, args.transcript_gap),
        )
        state.update({
            "words": [asdict(x) for x in words],
            "utterances": [asdict(x) for x in utterances],
            "subtitles": [asdict(x) for x in subtitles],
        })
        save_pipeline_state(state_path, state)
    else:
        print("[4/8] Skipping alignment, loading cached results")
        words, utterances, subtitles = load_alignment_artifacts(state)

    if step_enabled(args.start_from, "scenes"):
        print("[5/8] Detecting scenes / slides")
        scenes = detect_scene_boundaries(
            work_video,
            fps=fps,
            detector_name=args.scene_detector,
            threshold=args.scene_threshold,
            min_scene_sec=args.min_scene_sec,
        )
        scenes = merge_short_scenes(scenes, min_keep_sec=max(args.min_scene_sec, 6.0))

        if args.presentation == "yes":
            has_presentation = True
        elif args.presentation == "no":
            has_presentation = False
        else:
            has_presentation = decide_has_presentation(scenes, duration)

        print(f"[INFO] scenes={len(scenes)} has_presentation={has_presentation}")
        state.update({
            "scenes": [list(x) for x in scenes],
            "has_presentation": has_presentation,
        })
        save_pipeline_state(state_path, state)
    else:
        print("[5/8] Skipping scene detection, loading cached results")
        scenes, has_presentation = load_scene_artifacts(state)
        print(f"[INFO] scenes={len(scenes)} has_presentation={has_presentation}")

    pdf_generated = False
    if step_enabled(args.start_from, "slides"):
        print("[6/8] Building slide deck")
        slide_segments = build_slide_deck(
            scenes=scenes,
            utterances=utterances,
            work_video=work_video,
            out_dir=out_dir,
            force_single_slide=not has_presentation,
        )
        create_pptx(slide_segments, pptx_path, title=input_video.stem)

        if args.export_pdf:
            print("[6b/8] Converting PPTX to PDF")
            pdf_result = convert_pptx_to_pdf(pptx_path, out_dir)
            if pdf_result is None:
                print("[WARN] PDF conversion skipped: 'soffice' not found.")
            else:
                pdf_generated = True

        state["slide_segments"] = [asdict(x) for x in slide_segments]
        save_pipeline_state(state_path, state)
    else:
        print("[6/8] Skipping slide generation, loading cached results")
        slide_segments = load_slide_artifacts(state, pptx_path, args.export_pdf, pdf_path)
        pdf_generated = args.export_pdf and pdf_path.exists()

    if step_enabled(args.start_from, "write"):
        print("[7/8] Writing transcript and subtitle files")
        write_transcript_txt(utterances, transcript_txt)
        write_srt(subtitles, srt_path)
        write_ass(subtitles, ass_path, width=width, height=height)

        payload = {
            "input_video": str(input_video),
            "duration_sec": duration,
            "video": {"width": width, "height": height, "fps": fps},
            "asr_meta": asr_meta,
            "speaker_turns": [asdict(x) for x in turns],
            "words": [asdict(x) for x in words],
            "utterances": [asdict(x) for x in utterances],
            "subtitles": [asdict(x) for x in subtitles],
            "scenes": [list(x) for x in scenes],
            "has_presentation": has_presentation,
            "slide_segments": [asdict(x) for x in slide_segments],
        }
        save_json(payload, transcript_json)
        state["transcript_json"] = str(transcript_json)
        save_pipeline_state(state_path, state)
    else:
        print("[7/8] Skipping transcript/subtitle generation, reusing existing files")
        require_existing_file(srt_path, "subtitles.srt")
        require_existing_file(ass_path, "subtitles.ass")

    hard_ok = False
    if step_enabled(args.start_from, "render"):
        print("[8/8] Rendering output video")
        try:
            burn_subtitles_hard(work_video, ass_path, hard_mp4, video_encoder=video_encoder)
            hard_ok = True
        except subprocess.CalledProcessError as e:
            print("[WARN] Hard subtitles failed. Usually this means ffmpeg build has no libass.")
            if e.stderr:
                print(e.stderr)
        mux_subtitles_soft(work_video, srt_path, soft_mp4)
    else:
        print("[8/8] Skipping video rendering, reusing existing outputs")
        hard_ok = hard_mp4.exists()
        require_existing_file(soft_mp4, "video_softsubs.mp4")

    if not args.keep_work_files:
        shutil.rmtree(work_dir, ignore_errors=True)

    print("\nDone.")
    print(f"State JSON: {state_path}")
    if transcript_txt.exists():
        print(f"Transcript: {transcript_txt}")
    if transcript_json.exists():
        print(f"Transcript JSON: {transcript_json}")
    if pptx_path.exists():
        print(f"Slides PPTX: {pptx_path}")
    if pdf_generated or pdf_path.exists():
        print(f"Slides PDF: {pdf_path}")
    if hard_ok:
        print(f"Hard-sub video: {hard_mp4}")
    print(f"Soft-sub video: {soft_mp4}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
